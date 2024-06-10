import os
import re
import warnings
import numpy as np
import pandas as pd
import logging
from typing import Any

import fitz  # pdf lib # pip install --upgrade pymupdf
import docx  # pip install python-docx 
import ebooklib
from ebooklib import epub  # , fb2
from pptx import Presentation  # pip install python-pptx
from bs4 import BeautifulSoup

import torch
import torch.nn.functional as F
from transformers import pipeline
from sentence_transformers import SentenceTransformer

warnings.filterwarnings("ignore")
device = torch.device('mps' if torch.backends.mps.is_available() else 'cuda' if torch.backends else 'cpu')
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(name)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)


class FilePreprocessor:
    def __init__(
            self,
            summarizer_model_name: str,
            summarizer_model_max_length: int,
            st_model_name: str
    ):
        assert isinstance(summarizer_model_max_length, int)

        self.summarizer_model_name = summarizer_model_name
        self.summarizer_model_max_length = summarizer_model_max_length
        self.st_model_name = st_model_name

    @staticmethod
    def normalise_text(text: str) -> str:
        """
        Normalise the text in the form of a string, removing technical symbols and stripping leading and trailing whitespace
        :param text: input string to be normalised
        :return: a normalised string
        """
        return re.sub(r'\s+', ' ', text).strip()

    @staticmethod
    def get_first_n_words(text: str, n_words: int = 500) -> str:
        """
        Get first N words from a text.
        :param text: text to be cropped
        :param n_words: maximum number of words to be returned
        :return: cropped text
        """
        words = text.split()
        first_n_words = ' '.join(words[:n_words])
        return first_n_words

    @staticmethod
    def get_pdf_text(file_path: str) -> str:
        """
        Function parses text from pdf files.
        :param file_path: path till the file.
        :return: File content.
        """
        text = ""
        with fitz.open(file_path) as pdf:
            for page_num in range(pdf.page_count):
                page = pdf.load_page(page_num)
                text += page.get_text("text")
        return text

    @staticmethod
    def get_text_file_text(file_path: str) -> str:
        """
        Function that read files of the text formats and returns its content.
        :param file_path: path till the file. Accepted file formats are: txt, docx, epub, fb2.
        :return: File content.
        """
        file_extension = os.path.splitext(file_path)[1].lower()

        if file_extension in ['.txt', '.dmate']:
            with open(file_path, 'r') as file:
                return file.read()
        elif file_extension == '.docx':
            doc = docx.Document(file_path)
            return "\n".join([paragraph.text for paragraph in doc.paragraphs])
        elif file_extension == '.epub':
            book = epub.read_epub(file_path)
            return " ".join([item.get_body_content() for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT)])
        # elif file_extension == '.fb2':
        #     book = fb2.read_fb2(file_path)
        #     return " ".join([item.get_body_content() for item in book.get_items_of_type(ebooklib.ITEM_BODY)])

    @staticmethod
    def get_html_text(file_path: str) -> str:
        """
        Function parses text from html files.
        :param file_path: path till the file.
        :return: File content.
        """
        with open(file_path, 'r') as file:
            html_content = file.read()
            soup = BeautifulSoup(html_content, 'html.parser')
            # Remove script and style tags
            for script in soup(["script", "style"]):
                script.extract()
            # Get the text
            text = soup.get_text()
        return text

    @staticmethod
    def get_excel_text(file_path: str) -> str:
        """
        Function that read files of Excel/csv formats and returns its content.
        :param file_path: path till the file. Accepted file formats are: xls, xlsx, csv
        :return: File content
        """
        if file_path.endswith('.csv'):
            df = pd.read_csv(file_path)
        elif file_path.endswith('.xls') or file_path.endswith('.xlsx'):
            df = pd.read_excel(file_path)
        else:
            return "Unsupported file format"

        return df.to_string(index=False)

    @staticmethod
    def get_ppt_text(file_path: str) -> str:
        """
        Function parses text from .ppt (PowerPoint) files.
        :param file_path: path till the file.
        :return: File content.
        """
        presentation = Presentation(file_path)
        all_text = ""
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    all_text += shape.text + "\n"
        return all_text

    def _text_summarize(self, text: str) -> str:
        """
        Function text_summarize compresses the text.
        :param text: text to be summarized
        :param max_length: maximum length of the summarized text
        :param model_name: model name
        """
        if len(text) <= int(self.summarizer_model_max_length):
            logger.info("The text is shorter than set max_length. Skipping stage of summarization")
            return text
        summarizer = pipeline("summarization", model=self.summarizer_model_name)
        summary = summarizer(
            text, max_length=self.summarizer_model_max_length, min_length=10, do_sample=False
        )[0]['summary_text']
        logger.info(f"Compression is {(1 - len(summary) / len(text)) * 100:.2f}%")
        return summary

    def _text_to_embed(self, text: str) -> np.array:
        """
        Function text_to_embed converts the text to a numpy array Embedding.
        :param text: text to be embedded
        :param st_model_name: Transformer model name
        :return: embedding of text
        """
        model = SentenceTransformer(model_name_or_path=self.st_model_name)
        text_emb = model.encode(
            text,
            normalize_embeddings=True
        )
        return text_emb

    def __call__(self, file_path: str) -> np.array:
        """
        Function that receives file path with clearly stated extension and outputs its embedding
        :param file_path: path to the file embedding of which is to be received
        :return: file embedding
        """
        file_extension = os.path.splitext(file_path)[1].lower()
        file_name = file_path.split("/")[-1]

        if file_extension in ('.txt', '.docx', '.epub', '.fb2', '.dmate'):
            file_text = self.get_text_file_text(file_path)
        elif file_extension == '.pdf':
            file_text = self.get_pdf_text(file_path)
        elif file_extension in ('.csv', '.xls', '.xlsx'):
            file_text = self.get_excel_text(file_path)
        elif file_extension == '.ppt':
            file_text = self.get_ppt_text(file_path)
        elif file_extension == '.html':
            file_text = self.get_html_text(file_path)
        else:
            return "Unsupported file format"

        # compress text to ameliorate the performance of ML model
        raw_prompt = "File name: " + file_name + "file content: " + file_text  # add file name to the file text
        file_text_normalised = self.normalise_text(raw_prompt)  # remove noise from text
        file_text_cropped = self.get_first_n_words(file_text_normalised)  # cropping text to accelerate algorithms
        file_text_compressed = self._text_summarize(file_text_cropped)  # summarize text to get more info per word
        file_embed = self._text_to_embed(file_text_compressed)  # convert text to embedding
        return file_embed
